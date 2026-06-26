from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


BASE = Path(__file__).parent
ASSETS = BASE / "outputs" / "pitch_assets"
OUT = BASE / "outputs" / "pitch_assets" / "Churn_Pitch_3slides.pptx"
ASSETS.mkdir(parents=True, exist_ok=True)


def add_title(slide, text, subtitle=None):
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9.0), Inches(0.8))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(245, 245, 245)
    p.font.name = "Calibri"
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.52), Inches(1.0), Inches(9.0), Inches(0.45))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = RGBColor(186, 220, 255)
        sp.font.name = "Calibri"


def add_dark_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(13, 24, 45)


def add_bullets(slide, items, x, y, w, h, font_size=20):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(240, 245, 255)
        p.font.name = "Calibri"
        p.space_after = Pt(8)


def add_metric_card(slide, x, y, w, h, title, value):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(w), Inches(h)  # MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 52, 85)
    shape.line.color.rgb = RGBColor(56, 92, 140)

    t = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.12), Inches(w - 0.3), Inches(0.35))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(168, 205, 255)
    p.font.name = "Calibri"

    v = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.46), Inches(w - 0.3), Inches(0.5))
    vf = v.text_frame
    vf.clear()
    vp = vf.paragraphs[0]
    vp.text = value
    vp.font.size = Pt(24)
    vp.font.bold = True
    vp.font.color.rgb = RGBColor(255, 255, 255)
    vp.font.name = "Calibri"


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1 — Problem + headline results
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s1)
    add_title(
        s1,
        "Churn Intelligence: Кого удерживать и как",
        "Freedom 2026 Hackathon | Churn Prediction Track",
    )
    add_bullets(
        s1,
        [
            "• Мы решаем не только 'кто уйдет', а 'на кого стоит тратить retention-бюджет'.",
            "• Модель: CatBoost + поведенческие признаки + explainability (SHAP).",
            "• Выход: 4 action-сегмента для CRM (Persuadables / Sure Things / Sleeping Dogs / Lost Causes).",
        ],
        x=0.6,
        y=1.5,
        w=7.0,
        h=2.0,
        font_size=16,
    )

    add_metric_card(s1, 0.65, 4.0, 1.9, 1.2, "AUC-ROC", "0.8086")
    add_metric_card(s1, 2.75, 4.0, 1.9, 1.2, "F1-score", "0.6538")
    add_metric_card(s1, 4.85, 4.0, 1.9, 1.2, "Precision", "0.5555")
    add_metric_card(s1, 6.95, 4.0, 1.9, 1.2, "Recall", "0.7942")

    overview_img = ASSETS / "dashboard_live.png"
    if overview_img.exists():
        s1.shapes.add_picture(str(overview_img), Inches(8.2), Inches(1.35), Inches(4.7), Inches(5.6))

    # Slide 2 — How model learns + SHAP
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s2)
    add_title(s2, "Как обучалась модель и что она увидела")

    add_bullets(
        s2,
        [
            "• Данные: feature_store + events + transactions + acquisition + pLTV.",
            "• Label churn: нет активности 30+ дней.",
            "• Feature engineering: тренды активности, error/failed rates, product depth, digital ratio.",
            "• Контроль качества: дедупликация по user_id и clean merge без row explosion.",
        ],
        x=0.6,
        y=1.2,
        w=6.6,
        h=2.7,
        font_size=14,
    )

    add_bullets(
        s2,
        [
            "Ключевой инсайт:",
            "Поведение внутри приложения влияет сильнее, чем канал привлечения.",
            "is_paid_channel слабее, чем event/transaction динамика.",
        ],
        x=0.6,
        y=4.2,
        w=6.6,
        h=1.7,
        font_size=15,
    )

    shap_img = ASSETS / "feature_notes.png"
    if shap_img.exists():
        s2.shapes.add_picture(str(shap_img), Inches(7.4), Inches(1.25), Inches(5.5), Inches(5.7))

    # Slide 3 — Business actions + why solution is strong
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_dark_bg(s3)
    add_title(s3, "Результат для бизнеса: от модели к действиям")

    add_bullets(
        s3,
        [
            "Persuadables: персональный оффер/кэшбэк -> приоритет retention-бюджета.",
            "Sure Things: не тратим бюджет на удержание, делаем upsell.",
            "Sleeping Dogs: no-contact policy, чтобы не спровоцировать уход.",
            "Lost Causes: минимальный бюджет, только feedback loop.",
        ],
        x=0.6,
        y=1.25,
        w=6.6,
        h=2.6,
        font_size=14,
    )

    add_bullets(
        s3,
        [
            "Почему решение сильное:",
            "• Actionable, а не просто predictive.",
            "• Explainable (SHAP) -> понятно бизнесу.",
            "• Готово к запуску: dashboard + user-level scoring для CRM.",
        ],
        x=0.6,
        y=4.0,
        w=6.6,
        h=2.2,
        font_size=15,
    )

    top_users_img = ASSETS / "kpi_crop.png"
    eda_img = ASSETS / "dashboard_live.png"
    if top_users_img.exists():
        s3.shapes.add_picture(str(top_users_img), Inches(7.45), Inches(1.2), Inches(5.3), Inches(2.7))
    if eda_img.exists():
        s3.shapes.add_picture(str(eda_img), Inches(7.45), Inches(4.05), Inches(5.3), Inches(2.7))

    prs.save(OUT)
    print(f"Saved: {OUT}")


if __name__ == "__main__":
    main()

